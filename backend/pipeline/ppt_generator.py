"""
PowerPoint catalog generator.
Creates one slide per style group matching the reference PPT layout:
- Cover slide (title page)
- Style slides: front (left), back (center), detail (top right), specs (right panel)
- Logo, style name + number in header, company branding
"""

import os
import logging
from pathlib import Path
from typing import Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

from backend.models.schemas import StyleGroup, ImageJob, SpecData
from backend.utils.file_utils import get_catalog_output_path, get_reference_ppt, get_new_catalog_output_path

logger = logging.getLogger(__name__)

# Slide dimensions (widescreen 13.33 x 7.5 inches)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors matching the reference PPT (dark theme with cream accents)
DARK_BG = RGBColor(0x2D, 0x2A, 0x26)        # Dark brown/charcoal
HEADER_BG = RGBColor(0x2D, 0x2A, 0x26)       # Same dark for header
CREAM_TEXT = RGBColor(0xF5, 0xF0, 0xE8)       # Cream/off-white text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xE8, 0xE4, 0xDE)       # Light grey for image borders
DARK_TEXT = RGBColor(0x4A, 0x45, 0x3D)         # Dark text for specs
ACCENT_RED = RGBColor(0xC4, 0x3E, 0x3E)       # Asmara logo red
LINE_COLOR = RGBColor(0x8A, 0x84, 0x7C)       # Subtle line color
BODY_BG = RGBColor(0xF5, 0xF0, 0xE8)          # Cream body background


def _add_centered_picture(slide, img_path: str, box_left, box_top, box_width, box_height) -> None:
    """Helper to add an image to a slide while preserving aspect ratio and centering it within a box."""
    try:
        with Image.open(img_path) as img:
            img_w, img_h = img.size
    except Exception as e:
        logger.error(f"Failed to read image {img_path} for PPT size calculation: {e}")
        return

    # Calculate scaling to fit within the target box
    ratio = min(box_width / img_w, box_height / img_h)
    target_w = img_w * ratio
    target_h = img_h * ratio
    
    # Calculate centered position
    final_left = box_left + (box_width - target_w) / 2
    final_top = box_top + (box_height - target_h) / 2
    
    slide.shapes.add_picture(img_path, int(final_left), int(final_top), int(target_w), int(target_h))

def _add_framed_picture(slide, img_path: str, box_left, box_top, box_width, box_height, label_text: str = None) -> float:
    """Adds an image with a tight-fitting thin border frame. Returns the bottom position of the frame."""
    try:
        with Image.open(img_path) as img:
            img_w, img_h = img.size
    except Exception as e:
        logger.error(f"Failed to read image {img_path} for frame calculation: {e}")
        _add_centered_picture(slide, img_path, box_left, box_top, box_width, box_height)
        return box_top + box_height

    ratio = min(box_width / img_w, box_height / img_h)
    actual_w = int(img_w * ratio)
    actual_h = int(img_h * ratio)
    
    img_left = box_left + (box_width - actual_w) / 2
    img_top_pos = box_top + (box_height - actual_h) / 2

    border_pad = Inches(0.05)
    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        int(img_left - border_pad),
        int(img_top_pos - border_pad),
        int(actual_w + border_pad * 2),
        int(actual_h + border_pad * 2)
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = BODY_BG
    frame.line.color.rgb = LIGHT_GREY
    frame.line.width = Pt(1)

    slide.shapes.add_picture(img_path, int(img_left), int(img_top_pos), int(actual_w), int(actual_h))
    
    bottom_pos = img_top_pos + actual_h + border_pad

    # Add label below if provided
    if label_text:
        label_box = slide.shapes.add_textbox(
            int(img_left - border_pad), int(bottom_pos + Inches(0.05)),
            int(actual_w + border_pad * 2), Inches(0.25)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label_text
        p.font.size = Pt(9)
        p.font.color.rgb = DARK_TEXT
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        
        bottom_pos += Inches(0.3)

    return bottom_pos


def _add_cover_slide(prs: Presentation, total_styles: int) -> None:
    """Create the cover/title slide matching the reference layout."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Full dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Asmara logo text at top left (placeholder since we don't have the logo file)
    logo_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.4), Inches(2), Inches(0.6))
    tf = logo_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "asmara"
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_RED
    p.font.bold = True

    # Main title: ELEMENTS
    title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ELEMENTS"
    p.font.size = Pt(72)
    p.font.color.rgb = CREAM_TEXT
    p.font.bold = True
    p.font.name = "Georgia"

    # Subtitle: COLLECTION — SS26
    sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(6), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "COLLECTION — SS26"
    p.font.size = Pt(14)
    p.font.color.rgb = CREAM_TEXT
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.LEFT

    # Horizontal line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.9), Inches(4.5), Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE_COLOR
    line.line.fill.background()

    # Company info
    info_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.1), Inches(4), Inches(0.4))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ASMARA INDIA  ·  SUSTAINABLE FABRICS"
    p.font.size = Pt(10)
    p.font.color.rgb = LINE_COLOR
    p.font.name = "Arial"

    # Style count at bottom right
    count_box = slide.shapes.add_textbox(Inches(10.5), Inches(6.8), Inches(2.5), Inches(0.4))
    tf = count_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{total_styles} STYLES"
    p.font.size = Pt(10)
    p.font.color.rgb = LINE_COLOR
    p.alignment = PP_ALIGN.RIGHT


def _add_style_slide(
    prs: Presentation,
    group: StyleGroup,
    jobs: Dict[str, ImageJob],
    slide_number: int,
    total_styles: int
) -> None:
    """
    Create a style slide matching the reference PPT layout.
    Layout:
    - Header bar (dark): style number | style name + subtitle | logo
    - Body (cream): front image (left) | back image (center) | detail image + specs (right)
    - Footer: company name (left) | page number (right)
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # ── Background (Cover baked-in layout shapes) ──
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BODY_BG
    
    cover = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    cover.fill.solid()
    cover.fill.fore_color.rgb = BODY_BG
    cover.line.fill.background()

    # ── Header bar (dark strip at top) ──
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.85)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = HEADER_BG
    header.line.fill.background()

    # Style number (large, left side)
    num_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(0.8), Inches(0.7))
    tf = num_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"{slide_number:02d}"
    p.font.size = Pt(28)
    p.font.color.rgb = CREAM_TEXT
    p.font.bold = True
    p.font.name = "Georgia"

    # Vertical separator line
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(0.15), Pt(2), Inches(0.55)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = LINE_COLOR
    sep.line.fill.background()

    # Style name
    spec_data = group.spec_data or SpecData()
    style_title = group.name.upper() if group.name else f"STYLE {slide_number}"
    name_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.08), Inches(9.5), Inches(0.45))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = style_title
    p.font.size = Pt(22)
    p.font.color.rgb = CREAM_TEXT
    p.font.bold = True
    p.font.name = "Arial"

    # Subtitle: ref + fabric + GSM
    subtitle_parts = []
    if spec_data.ref_number:
        subtitle_parts.append(spec_data.ref_number)
    if spec_data.fabric_composition:
        subtitle_parts.append(spec_data.fabric_composition)
    if spec_data.gsm:
        gsm_val = spec_data.gsm.replace("g/m²", "").replace("g/m2", "").strip()
        subtitle_parts.append(f"{gsm_val} GSM")

    subtitle_text = "  ·  ".join(subtitle_parts) if subtitle_parts else ""
    sub_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.50), Inches(9.5), Inches(0.3))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle_text
    p.font.size = Pt(9)
    p.font.color.rgb = LINE_COLOR
    p.font.name = "Arial"


    # ── Image area ──
    img_top = Inches(1.1)
    img_height = Inches(5.0)

    # Front image (left)
    if group.front_image_id and group.front_image_id in jobs:
        front_job = jobs[group.front_image_id]
        img_path = front_job.processed_path or front_job.original_path
        if os.path.exists(img_path):
            _add_framed_picture(
                slide, img_path, Inches(0.5), img_top, Inches(3.8), img_height, label_text="F R O N T"
            )

    # Back image (center)
    if group.back_image_id and group.back_image_id in jobs:
        back_job = jobs[group.back_image_id]
        img_path = back_job.processed_path or back_job.original_path
        if os.path.exists(img_path):
            _add_framed_picture(
                slide, img_path, Inches(4.5), img_top, Inches(3.8), img_height, label_text="B A C K"
            )

    # Detail image (top right)
    detail_top = img_top
    detail_height = Inches(2.5)
    detail_bottom = detail_top + detail_height  # default if no detail image
    if group.detail_image_id and group.detail_image_id in jobs:
        detail_job = jobs[group.detail_image_id]
        img_path = detail_job.processed_path or detail_job.original_path
        if os.path.exists(img_path):
            detail_bottom = _add_framed_picture(
                slide, img_path, Inches(8.65), detail_top, Inches(4.2), detail_height
            )

    # ── Spec data panel (right side, below detail image) ──
    spec_top = detail_bottom + Inches(0.8)

    spec_fields = [
        ("REF", spec_data.ref_number),
        ("CONTENT", spec_data.fabric_composition),
        ("GSM", spec_data.gsm),
        ("REMARKS", spec_data.remarks),
        ("DATE", spec_data.date),
    ]

    y_offset = spec_top
    for label, value in spec_fields:
        if not value:
            continue

        # Label
        label_box = slide.shapes.add_textbox(
            Inches(8.7), y_offset, Inches(1.2), Inches(0.25)
        )
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.color.rgb = DARK_TEXT
        p.font.bold = True
        p.font.name = "Arial"

        # Value
        val_box = slide.shapes.add_textbox(
            Inches(10.0), y_offset, Inches(3.0), Inches(0.25)
        )
        tf = val_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(8)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Arial"

        y_offset += Inches(0.28)

    # ── Footer ──
    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.05), SLIDE_WIDTH, Inches(0.45)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = HEADER_BG
    footer.line.fill.background()

    # Company name (left)
    company_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(3), Inches(0.3))
    tf = company_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ASMARA INDIA"
    p.font.size = Pt(8)
    p.font.color.rgb = LINE_COLOR
    p.font.name = "Arial"

    # Page number (center)
    page_box = slide.shapes.add_textbox(Inches(5.66), Inches(7.1), Inches(2), Inches(0.3))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{slide_number} / {total_styles}"
    p.font.size = Pt(8)
    p.font.color.rgb = LINE_COLOR
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"

    # Logo at bottom right
    logo_box = slide.shapes.add_textbox(Inches(10.8), Inches(7.05), Inches(2.2), Inches(0.45))
    tf = logo_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "asmara"
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_RED
    p.font.bold = True
    p.alignment = PP_ALIGN.RIGHT


def generate_catalog(
    style_groups: Dict[str, StyleGroup],
    jobs: Dict[str, ImageJob]
) -> str:
    """
    Generate the complete catalog PowerPoint.
    Returns the output file path.
    """
    logger.info(f"Generating catalog with {len(style_groups)} style groups...")

    # Try to load reference PPT for dimensions and layout clues
    ref_ppt = get_reference_ppt()
    if ref_ppt.exists():
        logger.info(f"Loading reference PPT: {ref_ppt}")
        prs = Presentation(str(ref_ppt))
        # Clear existing slides by creating fresh presentation with same dimensions
        ref_width = prs.slide_width
        ref_height = prs.slide_height
        prs = Presentation()
        prs.slide_width = ref_width
        prs.slide_height = ref_height
        logger.info(f"Using reference dimensions: {ref_width}x{ref_height}")
    else:
        logger.info("No reference PPT found, using default dimensions")
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

    # Sort groups by style number
    sorted_groups = sorted(style_groups.values(), key=lambda g: g.style_number)

    # Auto-merge orphaned Spec Labels into the previous slide if it is missing spec data
    merged_groups = []
    for g in sorted_groups:
        has_images = g.front_image_id or g.back_image_id or g.detail_image_id
        if not has_images and g.spec_label_id and merged_groups:
            prev_g = merged_groups[-1]
            if not prev_g.spec_label_id:
                prev_g.spec_label_id = g.spec_label_id
                prev_g.spec_data = g.spec_data
                continue  # skip adding this empty group as a separate slide
        merged_groups.append(g)

    total_styles = len(merged_groups)

    # Cover slide
    _add_cover_slide(prs, total_styles)

    # Style slides
    for idx, group in enumerate(merged_groups, 1):
        _add_style_slide(prs, group, jobs, idx, total_styles)

    # Save
    output_path = str(get_new_catalog_output_path())
    prs.save(output_path)
    logger.info(f"Catalog saved to {output_path}")

    return output_path
